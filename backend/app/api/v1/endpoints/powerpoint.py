from fastapi import APIRouter, Depends, HTTPException, status
from sqlalchemy.ext.asyncio import AsyncSession
from sqlalchemy import select
from pydantic import BaseModel
from typing import List, Optional, Dict, Any
import uuid
from pathlib import Path
from datetime import datetime
import pandas as pd
import openpyxl
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from PIL import Image
import os
import shutil
import json

from app.core.database import get_db
from app.core.auth import get_current_user_optional
from app.models import Project, Run, Artifact, DrawingFolder, DrawingImage, AppUser
from app.core.storage import local_storage

router = APIRouter()

class ExcelSheetInfo(BaseModel):
    name: str
    columns: List[str]

class ExcelFileInfo(BaseModel):
    sheets: List[ExcelSheetInfo]

class RunImageInfo(BaseModel):
    filename: str
    path: str
    url: str

class DrawingImageInfo(BaseModel):
    id: str
    filename: str
    folder_id: str
    folder_description: Optional[str]
    url: str

class PowerPointConfig(BaseModel):
    run_id: str
    project_id: str
    excel_sheet: str = "meta"
    title_column: str = "main_level"
    description_column: Optional[str] = "description"
    match_column: str = "main_level"
    drawing_folder_id: Optional[str] = None
    layout: Dict[str, Any] = {}
    extra_text_columns: List[str] = []
    extra_image_folders: List[str] = []

async def check_project_access_for_powerpoint(
    db: AsyncSession, 
    project_id: uuid.UUID, 
    user: Optional[AppUser]
) -> Project:
    """Check if user has access to project - ensures members and owners can access."""
    result = await db.execute(select(Project).where(Project.id == project_id))
    project = result.scalar_one_or_none()
    
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")
    
    # Check if user is owner
    if user and project.owner_id == user.id:
        return project
    
    # Check if user is a member
    if user:
        from app.models import ProjectMember
        member_result = await db.execute(
            select(ProjectMember).where(
                ProjectMember.project_id == project_id,
                ProjectMember.user_id == user.id
            )
        )
        member = member_result.scalar_one_or_none()
        if member:
            return project
    
    # Allow guest access if project allows it
    if not user:
        if project.allow_guest:
            return project
        raise HTTPException(status_code=403, detail="Authentication required")
    
    # User is authenticated but not owner or member, and guest access is not allowed
    raise HTTPException(status_code=403, detail="Access denied - you are not a member of this project")

@router.get("/workspace/{workspace_id}/excel-info", response_model=ExcelFileInfo)
async def get_workspace_excel_info(
    workspace_id: str,
    db: AsyncSession = Depends(get_db),
    current_user: Optional[AppUser] = Depends(get_current_user_optional)
):
    """Get Excel file information from workspace."""
    
    workspace_dir = local_storage.base_path / "outputs" / "powerpoint" / workspace_id
    
    if not workspace_dir.exists():
        raise HTTPException(status_code=404, detail="Workspace not found")
    
    # Find Excel file in workspace
    excel_dir = workspace_dir / "excel"
    excel_path = None
    
    if excel_dir.exists():
        for ext in ['.xlsx', '.xls', '.xlsm']:
            excel_files = list(excel_dir.glob(f'*{ext}'))
            if excel_files:
                excel_path = excel_files[0]
                break
    
    if not excel_path or not excel_path.exists():
        raise HTTPException(status_code=404, detail="Excel file not found in workspace")
    
    try:
        # Load workbook
        wb = openpyxl.load_workbook(excel_path, data_only=True)
        
        sheets_info = []
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            # Get first row as column headers
            columns = []
            if sheet.max_row > 0:
                for cell in sheet[1]:
                    if cell.value:
                        columns.append(str(cell.value))
            
            sheets_info.append(ExcelSheetInfo(
                name=sheet_name,
                columns=columns
            ))
        
        return ExcelFileInfo(sheets=sheets_info)
    
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Failed to read Excel file: {str(e)}")

@router.get("/runs/{run_id}/excel-info", response_model=ExcelFileInfo)
async def get_run_excel_info(
    run_id: str,
    db: AsyncSession = Depends(get_db),
    current_user: Optional[AppUser] = Depends(get_current_user_optional)
):
    """Get Excel file information (sheets and columns) from a run."""
    
    # Get run and check access
    result = await db.execute(select(Run).where(Run.id == uuid.UUID(run_id)))
    run = result.scalar_one_or_none()
    
    if not run:
        raise HTTPException(status_code=404, detail="Run not found")
    
    await check_project_access_for_powerpoint(db, run.project_id, current_user)
    
    # Find Excel artifact
    artifacts_result = await db.execute(
        select(Artifact).where(
            Artifact.run_id == uuid.UUID(run_id),
            Artifact.kind == "input_excel"
        )
    )
    excel_artifact = artifacts_result.scalar_one_or_none()
    
    excel_path = None
    
    if excel_artifact:
        # Read Excel file - handle both absolute and relative paths
        storage_key = excel_artifact.storage_key
        
        # If storage_key is already an absolute path, use it directly
        if Path(storage_key).is_absolute():
            excel_path = Path(storage_key)
        else:
            # Otherwise use local_storage to resolve relative path
            excel_path = local_storage.get_file_path(storage_key)
    
    # If no artifact found or file doesn't exist, try to find Excel in run folder
    if not excel_path or not excel_path.exists():
        # Try to find Excel file in run folder
        run_dir_key = f"runs/{run_id}"
        run_dir_path = local_storage.get_file_path(run_dir_key)
        
        if run_dir_path.exists():
            # Look for Excel files in run folder
            excel_extensions = ['.xlsx', '.xls', '.xlsm']
            for ext in excel_extensions:
                for excel_file in run_dir_path.glob(f'*{ext}'):
                    if excel_file.exists():
                        excel_path = excel_file
                        break
                if excel_path:
                    break
    
    if not excel_path or not excel_path.exists():
        raise HTTPException(
            status_code=404, 
            detail=f"Excel file not found for this run. Please ensure the run has an Excel file (input_excel artifact or .xlsx/.xls file in run folder)."
        )
    
    try:
        # Load workbook
        wb = openpyxl.load_workbook(excel_path, data_only=True)
        
        sheets_info = []
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            # Get first row as column headers
            columns = []
            if sheet.max_row > 0:
                for cell in sheet[1]:
                    if cell.value:
                        columns.append(str(cell.value))
            
            sheets_info.append(ExcelSheetInfo(
                name=sheet_name,
                columns=columns
            ))
        
        return ExcelFileInfo(sheets=sheets_info)
    
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Failed to read Excel file: {str(e)}")

@router.get("/runs/{run_id}/images", response_model=List[RunImageInfo])
async def get_run_images(
    run_id: str,
    db: AsyncSession = Depends(get_db),
    current_user: Optional[AppUser] = Depends(get_current_user_optional)
):
    """Get all images from a run's task folder."""
    
    # Get run and check access
    result = await db.execute(select(Run).where(Run.id == uuid.UUID(run_id)))
    run = result.scalar_one_or_none()
    
    if not run:
        raise HTTPException(status_code=404, detail="Run not found")
    
    await check_project_access_for_powerpoint(db, run.project_id, current_user)
    
    # Find task folder - look for task_id in run
    if not run.jmp_task_id:
        return []
    
    from app.core.config import settings
    tasks_root = Path(settings.TASKS_DIRECTORY).expanduser().resolve()
    task_dir = tasks_root / f"task_{run.jmp_task_id}"
    
    if not task_dir.exists():
        return []
    
    # Find all image files in task directory
    image_extensions = {'.png', '.jpg', '.jpeg', '.gif', '.bmp', '.tiff'}
    images = []
    
    for file_path in task_dir.rglob('*'):
        if file_path.is_file() and file_path.suffix.lower() in image_extensions:
            relative_path = file_path.relative_to(tasks_root)
            # Create URL for serving the image
            import base64
            encoded_path = base64.b64encode(str(relative_path).encode()).decode()
            
            images.append(RunImageInfo(
                filename=file_path.name,
                path=str(relative_path),
                url=f"/api/v1/uploads/file-serve?path={encoded_path}"
            ))
    
    return sorted(images, key=lambda x: x.filename)

@router.get("/projects/{project_id}/drawing-folders/{folder_id}/images", response_model=List[DrawingImageInfo])
async def get_drawing_folder_images(
    project_id: str,
    folder_id: str,
    db: AsyncSession = Depends(get_db),
    current_user: Optional[AppUser] = Depends(get_current_user_optional)
):
    """Get all images from a drawing folder."""
    
    # Check project access
    await check_project_access_for_powerpoint(db, uuid.UUID(project_id), current_user)
    
    # Get folder
    folder_result = await db.execute(
        select(DrawingFolder).where(
            DrawingFolder.id == uuid.UUID(folder_id),
            DrawingFolder.project_id == uuid.UUID(project_id)
        )
    )
    folder = folder_result.scalar_one_or_none()
    
    if not folder:
        raise HTTPException(status_code=404, detail="Drawing folder not found")
    
    # Get images
    images_result = await db.execute(
        select(DrawingImage).where(DrawingImage.folder_id == uuid.UUID(folder_id))
    )
    images = images_result.scalars().all()
    
    image_infos = []
    backend_url = os.getenv("BACKEND_URL", "http://localhost:4700")
    
    for image in images:
        image_infos.append(DrawingImageInfo(
            id=str(image.id),
            filename=image.filename,
            folder_id=str(image.folder_id),
            folder_description=folder.description,
            url=f"{backend_url}/api/v1/projects/{project_id}/drawing-folders/{folder_id}/images/{image.id}/view"
        ))
    
    return sorted(image_infos, key=lambda x: x.filename)

class PrepareWorkspaceRequest(BaseModel):
    drawing_folder_ids: List[str] = []

@router.post("/runs/{run_id}/prepare-workspace")
async def prepare_workspace(
    run_id: str,
    request: PrepareWorkspaceRequest,
    db: AsyncSession = Depends(get_db),
    current_user: Optional[AppUser] = Depends(get_current_user_optional)
):
    """Prepare workspace by creating UUID folder and copying files."""
    import shutil
    
    # Get run and check access
    result = await db.execute(select(Run).where(Run.id == uuid.UUID(run_id)))
    run = result.scalar_one_or_none()
    
    if not run:
        raise HTTPException(status_code=404, detail="Run not found")
    
    await check_project_access_for_powerpoint(db, run.project_id, current_user)
    
    # Create UUID output folder
    workspace_id = str(uuid.uuid4())
    workspace_dir = local_storage.base_path / "outputs" / "powerpoint" / workspace_id
    workspace_dir.mkdir(parents=True, exist_ok=True)
    
    # Subdirectories
    excel_dir = workspace_dir / "excel"
    images_dir = workspace_dir / "images"
    excel_dir.mkdir(exist_ok=True)
    images_dir.mkdir(exist_ok=True)
    
    # Copy Excel file
    excel_path = None
    artifacts_result = await db.execute(
        select(Artifact).where(
            Artifact.run_id == uuid.UUID(run_id),
            Artifact.kind == "input_excel"
        )
    )
    excel_artifact = artifacts_result.scalar_one_or_none()
    
    if excel_artifact:
        storage_key = excel_artifact.storage_key
        if Path(storage_key).is_absolute():
            source_excel = Path(storage_key)
        else:
            source_excel = local_storage.get_file_path(storage_key)
        
        if source_excel.exists():
            dest_excel = excel_dir / source_excel.name
            shutil.copy2(source_excel, dest_excel)
            excel_path = str(dest_excel.relative_to(local_storage.base_path))
    
    # If no artifact, try to find Excel in run folder
    if not excel_path:
        run_dir_key = f"runs/{run_id}"
        run_dir_path = local_storage.get_file_path(run_dir_key)
        if run_dir_path.exists():
            for ext in ['.xlsx', '.xls', '.xlsm']:
                for excel_file in run_dir_path.glob(f'*{ext}'):
                    dest_excel = excel_dir / excel_file.name
                    shutil.copy2(excel_file, dest_excel)
                    excel_path = str(dest_excel.relative_to(local_storage.base_path))
                    break
                if excel_path:
                    break
    
    # Copy run task images
    from app.core.config import settings
    tasks_root = Path(settings.TASKS_DIRECTORY).expanduser().resolve()
    task_dir = tasks_root / f"task_{run.jmp_task_id}" if run.jmp_task_id else None
    
    image_files = []
    if task_dir and task_dir.exists():
        image_extensions = {'.png', '.jpg', '.jpeg', '.gif', '.bmp', '.tiff'}
        for file_path in task_dir.rglob('*'):
            if file_path.is_file() and file_path.suffix.lower() in image_extensions:
                dest_image = images_dir / file_path.name
                shutil.copy2(file_path, dest_image)
                image_files.append(file_path.name)
    
    # Copy drawing folder images into separate subfolders (drawings_1, drawings_2, etc.)
    drawing_folder_map = {}  # Map folder_id -> subfolder_name (drawings_1, drawings_2, etc.)
    drawing_files_by_folder = {}  # Map folder_id -> list of filenames
    
    for idx, folder_id_str in enumerate(request.drawing_folder_ids, start=1):
        try:
            folder_id = uuid.UUID(folder_id_str)
            subfolder_name = f"drawings_{idx}"
            drawings_dir = workspace_dir / subfolder_name
            drawings_dir.mkdir(exist_ok=True)
            
            # Store mapping
            drawing_folder_map[folder_id_str] = subfolder_name
            drawing_files_by_folder[folder_id_str] = []
            
            # Get images from this folder
            images_result = await db.execute(
                select(DrawingImage).where(DrawingImage.folder_id == folder_id)
            )
            for image in images_result.scalars().all():
                image_path = local_storage.get_file_path(image.storage_key)
                if image_path.exists():
                    dest_drawing = drawings_dir / image.filename
                    shutil.copy2(image_path, dest_drawing)
                    drawing_files_by_folder[folder_id_str].append(image.filename)
        except (ValueError, Exception) as e:
            # Skip invalid folder IDs
            continue
    
    # Store workspace config with folder mappings
    workspace_config = {
        "drawing_folder_ids": request.drawing_folder_ids,
        "drawing_folder_map": drawing_folder_map
    }
    config_file = workspace_dir / "config.json"
    with open(config_file, 'w') as f:
        json.dump(workspace_config, f)
    
    return {
        "workspace_id": workspace_id,
        "workspace_path": f"outputs/powerpoint/{workspace_id}",
        "excel_file": excel_path,
        "image_files": image_files,
        "drawing_files_by_folder": drawing_files_by_folder
    }

@router.post("/workspace/{workspace_id}/generate-powerpoint")
async def generate_powerpoint(
    workspace_id: str,
    config: PowerPointConfig,
    db: AsyncSession = Depends(get_db),
    current_user: Optional[AppUser] = Depends(get_current_user_optional)
):
    """Generate PowerPoint from prepared workspace."""
    
    workspace_dir = local_storage.base_path / "outputs" / "powerpoint" / workspace_id
    
    if not workspace_dir.exists():
        raise HTTPException(status_code=404, detail="Workspace not found")
    
    # Get run to check access
    result = await db.execute(select(Run).where(Run.id == uuid.UUID(config.run_id)))
    run = result.scalar_one_or_none()
    
    if not run:
        raise HTTPException(status_code=404, detail="Run not found")
    
    await check_project_access_for_powerpoint(db, run.project_id, current_user)
    
    # Read Excel from workspace
    excel_dir = workspace_dir / "excel"
    excel_path = None
    
    if excel_dir.exists():
        for ext in ['.xlsx', '.xls', '.xlsm']:
            excel_files = list(excel_dir.glob(f'*{ext}'))
            if excel_files:
                excel_path = excel_files[0]
                break
    
    if not excel_path or not excel_path.exists():
        raise HTTPException(status_code=404, detail="Excel file not found in workspace")
    
    try:
        # Read Excel sheet
        df = pd.read_excel(excel_path, sheet_name=config.excel_sheet)
        
        if config.title_column not in df.columns:
            raise HTTPException(status_code=400, detail=f"Title column '{config.title_column}' not found")
        
        if config.match_column not in df.columns:
            raise HTTPException(status_code=400, detail=f"Match column '{config.match_column}' not found")
        
        # Get images from workspace
        images_dir = workspace_dir / "images"
        
        run_images = {}
        if images_dir.exists():
            image_extensions = {'.png', '.jpg', '.jpeg', '.gif', '.bmp', '.tiff'}
            for file_path in images_dir.glob('*'):
                if file_path.is_file() and file_path.suffix.lower() in image_extensions:
                    filename_without_ext = file_path.stem
                    run_images[filename_without_ext] = str(file_path)
        
        # Load drawing images from multiple folders (drawings_1, drawings_2, etc.)
        drawing_images_by_folder = {}  # Map folder_id -> {filename_base: path}
        
        # Load workspace config to get folder mappings
        workspace_config_file = workspace_dir / "config.json"
        folder_map = {}
        if workspace_config_file.exists():
            with open(workspace_config_file, 'r') as f:
                workspace_config = json.load(f)
                folder_map = workspace_config.get('drawing_folder_map', {})
        
        # Load images from each drawings_N folder
        image_extensions = {'.png', '.jpg', '.jpeg', '.gif', '.bmp', '.tiff'}
        for folder_id_str, subfolder_name in folder_map.items():
            drawings_dir = workspace_dir / subfolder_name
            if drawings_dir.exists():
                drawing_images_by_folder[folder_id_str] = {}
                for file_path in drawings_dir.glob('*'):
                    if file_path.is_file() and file_path.suffix.lower() in image_extensions:
                        filename_without_ext = file_path.stem
                        drawing_images_by_folder[folder_id_str][filename_without_ext] = str(file_path)
        
        # For backward compatibility, create a combined drawing_images dict (uses first folder if multiple)
        drawing_images = {}
        if drawing_images_by_folder:
            # Use the first folder's images for the default drawingImage element
            first_folder_images = next(iter(drawing_images_by_folder.values()), {})
            drawing_images = first_folder_images
        
        # Create PowerPoint presentation
        prs = Presentation()
        
        # Default slide layout (blank)
        blank_slide_layout = prs.slide_layouts[6]
        
        # Track titles to skip duplicates
        seen_titles = set()
        
        # Group rows by title
        for _, row in df.iterrows():
            title = str(row[config.title_column]) if pd.notna(row[config.title_column]) else ""
            
            if not title or title in seen_titles:
                continue
            
            seen_titles.add(title)
            
            # Create slide
            slide = prs.slides.add_slide(blank_slide_layout)
            
            # Get match value for image matching
            match_value = str(row[config.match_column]) if pd.notna(row[config.match_column]) else ""
            # Normalize match value (remove any whitespace, convert to lowercase for matching)
            match_value_normalized = match_value.strip().lower()
            
            # Get layout from config if provided
            layout_elements = config.layout.get('elements', []) if isinstance(config.layout, dict) else []
            
            # Find elements from layout
            title_element = next((el for el in layout_elements if el.get('type') == 'title'), None)
            desc_element = next((el for el in layout_elements if el.get('type') == 'description'), None)
            run_image_element = next((el for el in layout_elements if el.get('type') == 'runImage'), None)
            drawing_image_element = next((el for el in layout_elements if el.get('type') == 'drawingImage'), None)
            
            # Find extra image elements that reference specific drawing folders
            extra_image_elements = [el for el in layout_elements if el.get('type') == 'extraImage' and el.get('folderId')]
            
            # Find extra text elements
            extra_text_elements = [el for el in layout_elements if el.get('type') == 'text' and el.get('column')]
            
            def hex_to_rgb(hex_color: str) -> tuple:
                """Convert hex color string to RGB tuple."""
                if not hex_color or not isinstance(hex_color, str):
                    return (0, 0, 0)  # Default to black
                hex_color = hex_color.lstrip('#')
                if len(hex_color) == 6:
                    try:
                        return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))
                    except (ValueError, IndexError):
                        return (0, 0, 0)  # Default to black on error
                return (0, 0, 0)  # Default to black
            
            # Get title element properties with defaults
            title_x = Inches(title_element.get('x', 0.5) if title_element else 0.5)
            title_y = Inches(title_element.get('y', 0.5) if title_element else 0.5)
            title_width = Inches(title_element.get('width', 9) if title_element else 9)
            title_height = Inches(title_element.get('height', 0.8) if title_element else 0.8)
            title_font_size = Pt(title_element.get('fontSize', 24) if title_element else 24)
            
            # Add title in separate text box
            title_box = slide.shapes.add_textbox(title_x, title_y, title_width, title_height)
            title_frame = title_box.text_frame
            title_frame.text = title
            title_frame.paragraphs[0].font.size = title_font_size
            title_frame.paragraphs[0].font.bold = title_element.get('bold', True) if title_element else True
            if title_element and title_element.get('color'):
                try:
                    color_str = title_element.get('color')
                    if color_str:
                        rgb = hex_to_rgb(color_str)
                        # Validate RGB values are in range [0, 255]
                        rgb = tuple(max(0, min(255, int(c))) for c in rgb)
                        title_frame.paragraphs[0].font.color.rgb = RGBColor(rgb[0], rgb[1], rgb[2])
                except (ValueError, TypeError, AttributeError):
                    pass  # Keep default color if conversion fails
            title_frame.word_wrap = True
            
            # Add description in separate text box if configured
            if config.description_column and config.description_column in df.columns:
                desc_value = str(row[config.description_column]) if pd.notna(row[config.description_column]) else ""
                if desc_value:
                    desc_x = Inches(desc_element.get('x', 0.5) if desc_element else 0.5)
                    desc_y = Inches(desc_element.get('y', 1.3) if desc_element else 1.3)
                    desc_width = Inches(desc_element.get('width', 9) if desc_element else 9)
                    desc_height = Inches(desc_element.get('height', 0.4) if desc_element else 0.4)
                    desc_font_size = Pt(desc_element.get('fontSize', 14) if desc_element else 14)
                    
                    desc_box = slide.shapes.add_textbox(desc_x, desc_y, desc_width, desc_height)
                    desc_frame = desc_box.text_frame
                    desc_frame.text = desc_value
                    desc_frame.paragraphs[0].font.size = desc_font_size
                    if desc_element and desc_element.get('color'):
                        try:
                            color_str = desc_element.get('color')
                            if color_str:
                                rgb = hex_to_rgb(color_str)
                                # Validate RGB values are in range [0, 255]
                                rgb = tuple(max(0, min(255, int(c))) for c in rgb)
                                desc_frame.paragraphs[0].font.color.rgb = RGBColor(rgb[0], rgb[1], rgb[2])
                        except (ValueError, TypeError, AttributeError) as e:
                            pass  # Keep default color if conversion fails
                    desc_frame.word_wrap = True
            
            # Get maintainAspectRatio settings (default True for images)
            maintain_run_aspect = run_image_element.get('maintainAspectRatio', True) if run_image_element else True
            maintain_drawing_aspect = drawing_image_element.get('maintainAspectRatio', True) if drawing_image_element else True
            
            # Default positions if not in layout
            left_x = Inches(run_image_element.get('x', 0.5) if run_image_element else 0.5)
            left_y = Inches(run_image_element.get('y', 1.5) if run_image_element else 1.5)
            left_width = Inches(run_image_element.get('width', 6) if run_image_element else 6)
            left_height = Inches(run_image_element.get('height', 5) if run_image_element else 5)
            
            right_x = Inches(drawing_image_element.get('x', 7) if drawing_image_element else 7)
            right_y = Inches(drawing_image_element.get('y', 1.5) if drawing_image_element else 1.5)
            right_width = Inches(drawing_image_element.get('width', 2.5) if drawing_image_element else 2.5)
            right_height = Inches(drawing_image_element.get('height', 5) if drawing_image_element else 5)
            
            # Match run images by exact filename (excluding extension)
            matched_run_image_path = None
            for filename_base, image_path in run_images.items():
                # Exact match (case-insensitive, excluding extension)
                if filename_base.lower() == match_value_normalized:
                    matched_run_image_path = image_path
                    break
            
            if matched_run_image_path:
                # Add matched run image
                if maintain_run_aspect:
                    # Preserve aspect ratio - calculate height based on width
                    try:
                        # Get original image dimensions to calculate aspect ratio
                        with Image.open(matched_run_image_path) as img:
                            original_width, original_height = img.size
                            aspect_ratio = original_height / original_width if original_width > 0 else 1
                        
                        # Calculate height based on width and aspect ratio
                        calculated_height = left_width * aspect_ratio
                        
                        # Use calculated height (but don't exceed the original height if specified in layout)
                        image_height = min(calculated_height, left_height) if left_height > 0 else calculated_height
                        
                        slide.shapes.add_picture(matched_run_image_path, left_x, left_y, left_width, image_height)
                    except Exception as e:
                        # Fallback to original dimensions if image can't be read
                        slide.shapes.add_picture(matched_run_image_path, left_x, left_y, left_width, left_height)
                else:
                    # Use exact dimensions from layout
                    slide.shapes.add_picture(matched_run_image_path, left_x, left_y, left_width, left_height)
            else:
                # Add "not found" note for run image
                not_found_box = slide.shapes.add_textbox(left_x, left_y, left_width, left_height)
                not_found_frame = not_found_box.text_frame
                not_found_frame.text = "Run Image\nNot Found"
                not_found_frame.paragraphs[0].font.size = Pt(16)
                # Don't set color - will default to black
                # Center align
                not_found_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            
            # Match drawing images by exact filename (excluding extension)
            matched_drawing_image_path = None
            if drawing_images:  # Only if drawing folder was selected
                for filename_base, image_path in drawing_images.items():
                    # Exact match (case-insensitive, excluding extension)
                    if filename_base.lower() == match_value_normalized:
                        matched_drawing_image_path = image_path
                        break
            
            if matched_drawing_image_path:
                # Add matched drawing image
                if maintain_drawing_aspect:
                    # Preserve aspect ratio - calculate height based on width
                    try:
                        # Get original image dimensions to calculate aspect ratio
                        with Image.open(matched_drawing_image_path) as img:
                            original_width, original_height = img.size
                            aspect_ratio = original_height / original_width if original_width > 0 else 1
                        
                        # Calculate height based on width and aspect ratio
                        calculated_height = right_width * aspect_ratio
                        
                        # Use calculated height (but don't exceed the original height if specified in layout)
                        image_height = min(calculated_height, right_height) if right_height > 0 else calculated_height
                        
                        slide.shapes.add_picture(matched_drawing_image_path, right_x, right_y, right_width, image_height)
                    except Exception as e:
                        # Fallback to original dimensions if image can't be read
                        slide.shapes.add_picture(matched_drawing_image_path, right_x, right_y, right_width, right_height)
                else:
                    # Use exact dimensions from layout
                    slide.shapes.add_picture(matched_drawing_image_path, right_x, right_y, right_width, right_height)
            elif drawing_images:  # Only show "not found" if drawing folder was selected
                # Add "not found" note for drawing image
                not_found_box = slide.shapes.add_textbox(right_x, right_y, right_width, right_height)
                not_found_frame = not_found_box.text_frame
                not_found_frame.text = "Drawing Image\nNot Found"
                not_found_frame.paragraphs[0].font.size = Pt(14)
                # Don't set color - will default to black
                # Center align
                not_found_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            
            # Add extra image elements from multiple drawing folders
            for extra_image_element in extra_image_elements:
                folder_id_str = extra_image_element.get('folderId')
                if folder_id_str and folder_id_str in drawing_images_by_folder:
                    # Get images from this specific folder
                    folder_images = drawing_images_by_folder[folder_id_str]
                    
                    # Match image from this folder
                    matched_extra_image_path = None
                    for filename_base, image_path in folder_images.items():
                        if filename_base.lower() == match_value_normalized:
                            matched_extra_image_path = image_path
                            break
                    
                    # Get element properties
                    extra_x = Inches(extra_image_element.get('x', 0.5))
                    extra_y = Inches(extra_image_element.get('y', 2.0))
                    extra_width = Inches(extra_image_element.get('width', 2))
                    extra_height = Inches(extra_image_element.get('height', 2))
                    maintain_extra_aspect = extra_image_element.get('maintainAspectRatio', True)
                    
                    if matched_extra_image_path:
                        # Add matched extra image
                        if maintain_extra_aspect:
                            try:
                                with Image.open(matched_extra_image_path) as img:
                                    original_width, original_height = img.size
                                    aspect_ratio = original_height / original_width if original_width > 0 else 1
                                    calculated_height = extra_width * aspect_ratio
                                    image_height = min(calculated_height, extra_height) if extra_height > 0 else calculated_height
                                    slide.shapes.add_picture(matched_extra_image_path, extra_x, extra_y, extra_width, image_height)
                            except Exception:
                                slide.shapes.add_picture(matched_extra_image_path, extra_x, extra_y, extra_width, extra_height)
                        else:
                            slide.shapes.add_picture(matched_extra_image_path, extra_x, extra_y, extra_width, extra_height)
                    else:
                        # Add "not found" note
                        not_found_box = slide.shapes.add_textbox(extra_x, extra_y, extra_width, extra_height)
                        not_found_frame = not_found_box.text_frame
                        not_found_frame.text = f"Image\nNot Found"
                        not_found_frame.paragraphs[0].font.size = Pt(12)
                        not_found_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            
            # Add extra text elements
            for text_element in extra_text_elements:
                column_name = text_element.get('column')
                if column_name and column_name in df.columns:
                    text_value = str(row[column_name]) if pd.notna(row[column_name]) else ""
                    if text_value:
                        text_x = Inches(text_element.get('x', 0.5))
                        text_y = Inches(text_element.get('y', 2.0))
                        text_width = Inches(text_element.get('width', 4))
                        text_height = Inches(text_element.get('height', 0.5))
                        text_font_size = Pt(text_element.get('fontSize', 12))
                        
                        text_box = slide.shapes.add_textbox(text_x, text_y, text_width, text_height)
                        text_frame = text_box.text_frame
                        text_frame.text = text_value
                        text_frame.paragraphs[0].font.size = text_font_size
                        if text_element.get('color'):
                            try:
                                color_str = text_element.get('color')
                                if color_str:
                                    rgb = hex_to_rgb(color_str)
                                    # Validate RGB values are in range [0, 255]
                                    rgb = tuple(max(0, min(255, int(c))) for c in rgb)
                                    text_frame.paragraphs[0].font.color.rgb = RGBColor(rgb[0], rgb[1], rgb[2])
                            except (ValueError, TypeError, AttributeError) as e:
                                pass  # Keep default color if conversion fails
                        if text_element.get('bold'):
                            text_frame.paragraphs[0].font.bold = True
                        text_frame.word_wrap = True
        
        # Generate unique output ID
        output_id = str(uuid.uuid4())
        pptx_filename = f"presentation_{output_id}.pptx"
        pptx_path = workspace_dir / pptx_filename
        
        prs.save(str(pptx_path))
        
        # Save settings JSON
        settings_data = {
            "output_id": output_id,
            "run_id": config.run_id,
            "project_id": config.project_id,
            "workspace_id": workspace_id,
            "excel_sheet": config.excel_sheet,
            "title_column": config.title_column,
            "description_column": config.description_column,
            "match_column": config.match_column,
            "layout": config.layout,
            "extra_text_columns": config.extra_text_columns,
            "extra_image_folders": config.extra_image_folders,
            "created_at": datetime.now().isoformat(),
            "slide_count": len(seen_titles)
        }
        settings_filename = f"settings_{output_id}.json"
        settings_path = workspace_dir / settings_filename
        with open(settings_path, 'w') as f:
            json.dump(settings_data, f, indent=2)
        
        # Store output metadata in run's outputs directory
        run_outputs_dir = local_storage.base_path / "outputs" / "powerpoint" / "runs" / config.run_id
        run_outputs_dir.mkdir(parents=True, exist_ok=True)
        output_metadata = {
            "output_id": output_id,
            "run_id": config.run_id,
            "workspace_id": workspace_id,
            "filename": pptx_filename,
            "settings_filename": settings_filename,
            "created_at": settings_data["created_at"],
            "slide_count": len(seen_titles)
        }
        output_metadata_path = run_outputs_dir / f"output_{output_id}.json"
        with open(output_metadata_path, 'w') as f:
            json.dump(output_metadata, f, indent=2)
        
        # Create download URL - use dedicated endpoint with access control
        return {
            "success": True,
            "output_id": output_id,
            "filename": pptx_filename,
            "download_url": f"/api/v1/powerpoint/outputs/{output_id}/download",
            "slide_count": len(seen_titles)
        }
    
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Failed to generate PowerPoint: {str(e)}")

class PowerPointOutputInfo(BaseModel):
    output_id: str
    run_id: str
    workspace_id: str
    filename: str
    created_at: str
    slide_count: int

@router.get("/runs/{run_id}/outputs", response_model=List[PowerPointOutputInfo])
async def get_run_outputs(
    run_id: str,
    db: AsyncSession = Depends(get_db),
    current_user: Optional[AppUser] = Depends(get_current_user_optional)
):
    """Get all PowerPoint outputs for a run."""
    
    # Get run and check access
    result = await db.execute(select(Run).where(Run.id == uuid.UUID(run_id)))
    run = result.scalar_one_or_none()
    
    if not run:
        raise HTTPException(status_code=404, detail="Run not found")
    
    await check_project_access_for_powerpoint(db, run.project_id, current_user)
    
    # List all outputs for this run
    run_outputs_dir = local_storage.base_path / "outputs" / "powerpoint" / "runs" / run_id
    
    outputs = []
    if run_outputs_dir.exists():
        for metadata_file in run_outputs_dir.glob("output_*.json"):
            try:
                with open(metadata_file, 'r') as f:
                    metadata = json.load(f)
                    outputs.append(PowerPointOutputInfo(
                        output_id=metadata.get("output_id", ""),
                        run_id=metadata.get("run_id", run_id),
                        workspace_id=metadata.get("workspace_id", ""),
                        filename=metadata.get("filename", ""),
                        created_at=metadata.get("created_at", ""),
                        slide_count=metadata.get("slide_count", 0)
                    ))
            except Exception:
                continue
    
    # Sort by created_at descending
    outputs.sort(key=lambda x: x.created_at, reverse=True)
    
    return outputs

@router.get("/outputs/{output_id}/settings")
async def get_output_settings(
    output_id: str,
    db: AsyncSession = Depends(get_db),
    current_user: Optional[AppUser] = Depends(get_current_user_optional)
):
    """Get settings for a specific PowerPoint output."""
    
    # Find output metadata
    outputs_base_dir = local_storage.base_path / "outputs" / "powerpoint" / "runs"
    output_metadata = None
    run_id = None
    
    # Search all runs for this output
    for run_dir in outputs_base_dir.iterdir():
        if run_dir.is_dir():
            metadata_file = run_dir / f"output_{output_id}.json"
            if metadata_file.exists():
                try:
                    with open(metadata_file, 'r') as f:
                        output_metadata = json.load(f)
                        run_id = output_metadata.get("run_id")
                        break
                except Exception:
                    continue
    
    if not output_metadata or not run_id:
        raise HTTPException(status_code=404, detail="Output not found")
    
    # Check access
    result = await db.execute(select(Run).where(Run.id == uuid.UUID(run_id)))
    run = result.scalar_one_or_none()
    
    if not run:
        raise HTTPException(status_code=404, detail="Run not found")
    
    await check_project_access_for_powerpoint(db, run.project_id, current_user)
    
    # Load settings from workspace
    workspace_id = output_metadata.get("workspace_id")
    settings_filename = output_metadata.get("settings_filename")
    
    if not workspace_id or not settings_filename:
        raise HTTPException(status_code=404, detail="Settings not found")
    
    workspace_dir = local_storage.base_path / "outputs" / "powerpoint" / workspace_id
    settings_path = workspace_dir / settings_filename
    
    if not settings_path.exists():
        raise HTTPException(status_code=404, detail="Settings file not found")
    
    try:
        with open(settings_path, 'r') as f:
            settings = json.load(f)
        return settings
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Failed to load settings: {str(e)}")

@router.delete("/outputs/{output_id}")
async def delete_powerpoint_output(
    output_id: str,
    db: AsyncSession = Depends(get_db),
    current_user: Optional[AppUser] = Depends(get_current_user_optional)
):
    """Delete a PowerPoint output - only project owners can delete."""
    import shutil
    
    # Find output metadata
    outputs_base_dir = local_storage.base_path / "outputs" / "powerpoint" / "runs"
    output_metadata = None
    run_id = None
    
    # Search all runs for this output
    for run_dir in outputs_base_dir.iterdir():
        if run_dir.is_dir():
            metadata_file = run_dir / f"output_{output_id}.json"
            if metadata_file.exists():
                try:
                    with open(metadata_file, 'r') as f:
                        output_metadata = json.load(f)
                        run_id = output_metadata.get("run_id")
                        break
                except Exception:
                    continue
    
    if not output_metadata or not run_id:
        raise HTTPException(status_code=404, detail="Output not found")
    
    # Check access via run's project
    result = await db.execute(select(Run).where(Run.id == uuid.UUID(run_id)))
    run = result.scalar_one_or_none()
    
    if not run:
        raise HTTPException(status_code=404, detail="Run not found")
    
    # Get project
    project_result = await db.execute(select(Project).where(Project.id == run.project_id))
    project = project_result.scalar_one_or_none()
    
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")
    
    # Only owners can delete
    if not current_user:
        raise HTTPException(status_code=401, detail="Authentication required")
    
    if project.owner_id != current_user.id:
        raise HTTPException(status_code=403, detail="Only project owners can delete presentations")
    
    # Get file paths
    workspace_id = output_metadata.get("workspace_id")
    filename = output_metadata.get("filename")
    settings_filename = output_metadata.get("settings_filename")
    
    # Delete metadata file
    metadata_file_path = outputs_base_dir / run_id / f"output_{output_id}.json"
    if metadata_file_path.exists():
        metadata_file_path.unlink()
    
    # Delete PowerPoint file
    if workspace_id and filename:
        pptx_file_path = local_storage.base_path / "outputs" / "powerpoint" / workspace_id / filename
        if pptx_file_path.exists():
            pptx_file_path.unlink()
    
    # Delete settings file
    if workspace_id and settings_filename:
        settings_file_path = local_storage.base_path / "outputs" / "powerpoint" / workspace_id / settings_filename
        if settings_file_path.exists():
            settings_file_path.unlink()
    
    return {"success": True, "message": "Presentation deleted successfully"}

@router.get("/outputs/{output_id}/download")
async def download_powerpoint_output(
    output_id: str,
    db: AsyncSession = Depends(get_db),
    current_user: Optional[AppUser] = Depends(get_current_user_optional)
):
    """Download a PowerPoint output file with project access check."""
    from fastapi.responses import FileResponse
    
    # Find output metadata
    outputs_base_dir = local_storage.base_path / "outputs" / "powerpoint" / "runs"
    output_metadata = None
    run_id = None
    
    # Search all runs for this output
    for run_dir in outputs_base_dir.iterdir():
        if run_dir.is_dir():
            metadata_file = run_dir / f"output_{output_id}.json"
            if metadata_file.exists():
                try:
                    with open(metadata_file, 'r') as f:
                        output_metadata = json.load(f)
                        run_id = output_metadata.get("run_id")
                        break
                except Exception:
                    continue
    
    if not output_metadata or not run_id:
        raise HTTPException(status_code=404, detail="Output not found")
    
    # Check access via run's project
    result = await db.execute(select(Run).where(Run.id == uuid.UUID(run_id)))
    run = result.scalar_one_or_none()
    
    if not run:
        raise HTTPException(status_code=404, detail="Run not found")
    
    # Check project access - this ensures only project members/owners can access
    await check_project_access_for_powerpoint(db, run.project_id, current_user)
    
    # Get file path from metadata
    workspace_id = output_metadata.get("workspace_id")
    filename = output_metadata.get("filename")
    
    if not workspace_id or not filename:
        raise HTTPException(status_code=404, detail="File information not found")
    
    # Construct file path
    file_path = local_storage.base_path / "outputs" / "powerpoint" / workspace_id / filename
    
    if not file_path.exists() or not file_path.is_file():
        raise HTTPException(status_code=404, detail="PowerPoint file not found")
    
    # Return file with proper MIME type
    return FileResponse(
        path=str(file_path),
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        filename=filename,
        headers={
            "Content-Disposition": f'attachment; filename="{filename}"'
        }
    )

